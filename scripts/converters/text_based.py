"""Converters for text/structured formats - no LLM needed."""
from pathlib import Path
from docx import Document
from pptx import Presentation
import mammoth


def convert_docx(path):
    """Try mammoth (better formatting), fallback to python-docx."""
    try:
        with open(path, "rb") as f:
            result = mammoth.convert_to_markdown(f)
        return result.value
    except Exception:
        doc = Document(path)
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def convert_pptx(path):
    """Extract text from each slide. Image-heavy slides should also be sent to Gemini for visual context."""
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slide_parts.append(f"**Notes:** {slide.notes_slide.notes_text_frame.text}")
        parts.append("\n\n".join(slide_parts))
    return "\n\n".join(parts)
